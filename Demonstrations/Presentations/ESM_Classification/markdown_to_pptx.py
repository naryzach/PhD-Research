import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def markdown_to_pptx(input_file="slides.md", output_file="ESM_Classification_Skeleton.pptx"):
    print(f"Parsing {input_file}...")
    
    if not os.path.exists(input_file):
        print(f"Error: {input_file} not found.")
        return

    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()

    # Split into slides
    slides_raw = re.split(r'^---', content, flags=re.MULTILINE)
    
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        # Strip frontmatter (between first --- and second --- or if at start)
        lines = slide_text.split('\n')
        clean_lines = []
        in_frontmatter = False
        for line in lines:
            if line.strip() == '---':
                in_frontmatter = not in_frontmatter
                continue
            if in_frontmatter:
                continue
            # Also skip style tags
            if '<style>' in line or '</style>' in line:
                continue
            clean_lines.append(line)
        
        slide_content = '\n'.join(clean_lines).strip()
        
        # Extract title (first line starting with #)
        title_match = re.search(r'^#\s+(.*)', slide_content, re.MULTILINE)
        title = title_match.group(1) if title_match else ""
        
        # Extract sub-title or body
        # Remove the title line from content
        body_content = re.sub(r'^#\s+.*', '', slide_content, count=1, flags=re.MULTILINE).strip()
        
        # Remove HTML tags for the text version
        body_content = re.sub(r'<[^>]+>', '', body_content)
        
        # Add slide
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        slide.shapes.title.text = title
        
        # Set Body
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        
        # Process lines to handle bullets
        body_lines = body_content.split('\n')
        for line in body_lines:
            line = line.strip()
            if not line: continue
            
            p = tf.add_paragraph()
            if line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
                p.text = line[2:]
                p.level = 0
            elif line.startswith('### '):
                p.text = line[4:]
                p.font.bold = True
                p.font.size = Pt(24)
            else:
                p.text = line
                p.level = 0

    prs.save(output_file)
    print(f"Exported text skeleton to: {output_file}")

if __name__ == "__main__":
    markdown_to_pptx()
