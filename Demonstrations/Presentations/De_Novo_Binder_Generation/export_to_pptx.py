import asyncio
import os
import re
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

async def export_slidev_to_pptx(url="http://localhost:3030", output_file="De_Novo_Binder_Generation.pptx"):
    print(f"🚀 Starting export from {url}...")
    
    # 1. Determine slide count from slides.md
    slide_count = 1
    if os.path.exists("slides.md"):
        with open("slides.md", "r") as f:
            content = f.read()
            # Slidev uses '---' as separators
            slide_count = len(re.findall(r'^---', content, re.MULTILINE)) + 1
    
    print(f"📊 Detected {slide_count} slides in slides.md")

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        # Set viewport to standard Slidev aspect ratio (16:9)
        page = await browser.new_page(viewport={'width': 1920, 'height': 1080})
        
        prs = Presentation()
        # Set PPTX slide size to 16:9 (13.333 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i in range(1, slide_count + 1):
            slide_url = f"{url}/{i}"
            print(f"📸 Capturing Slide {i}/{slide_count}...")
            
            try:
                await page.goto(slide_url, wait_until="networkidle")
                # Wait for components/animations to settle
                await asyncio.sleep(3) 
                
                screenshot_path = f"temp_slide_{i}.png"
                await page.screenshot(path=screenshot_path)
                
                # Add to PPTX
                slide_layout = prs.slide_layouts[6] # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Stretch image to fill the slide
                slide.shapes.add_picture(screenshot_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # Clean up temp image
                os.remove(screenshot_path)
            except Exception as e:
                print(f"❌ Error on slide {i}: {e}")

        prs.save(output_file)
        await browser.close()
        print(f"✅ Export complete! Saved to: {output_file}")

if __name__ == "__main__":
    asyncio.run(export_slidev_to_pptx())
